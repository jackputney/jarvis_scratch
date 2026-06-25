"""tools/pitch_deck.py — Generate pitch decks using Claude + python-pptx."""

from __future__ import annotations

import json
import os
import platform
import subprocess
from datetime import datetime
from pathlib import Path


def create_pitch_deck(
    topic: str,
    slide_count: int = 5,
    style: str = "professional",
    output_dir: str = "",
) -> str:
    """Generate a pitch deck .pptx file on the topic with slide_count slides."""
    topic = (topic or "").strip()
    if not topic:
        return "Refused: topic is required."

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ImportError:
        return "python-pptx is not installed. Run: pip install python-pptx"

    try:
        slide_count = max(1, min(int(slide_count), 15))
    except (TypeError, ValueError):
        slide_count = 5

    slides_data = _generate_slide_content(topic, slide_count, style)
    if not slides_data:
        return "Failed to generate slide content."

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    schemes = {
        "professional": {
            "bg": RGBColor(0x1A, 0x1E, 0x28),
            "accent": RGBColor(0x6C, 0x7B, 0xF7),
            "text": RGBColor(0xE8, 0xE9, 0xED),
        },
        "light": {
            "bg": RGBColor(0xFF, 0xFF, 0xFF),
            "accent": RGBColor(0x1A, 0x1E, 0x28),
            "text": RGBColor(0x1A, 0x1E, 0x28),
        },
        "bold": {
            "bg": RGBColor(0x00, 0x00, 0x00),
            "accent": RGBColor(0xFF, 0xD7, 0x00),
            "text": RGBColor(0xFF, 0xFF, 0xFF),
        },
    }
    colors = schemes.get((style or "").strip().lower(), schemes["professional"])
    blank_layout = prs.slide_layouts[6]

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["bg"]

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors["accent"]
        bar.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.8), Inches(0.4))
        num_tf = num_box.text_frame
        num_tf.text = f"{i + 1}/{len(slides_data)}"
        num_tf.paragraphs[0].runs[0].font.size = Pt(10)
        num_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x5A, 0x5F, 0x73)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.5), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(36) if i == 0 else Pt(28)
        p.font.bold = True
        p.font.color.rgb = colors["text"]

        if i == 0 and slide_data.get("subtitle"):
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(12.5), Inches(0.8))
            sub_tf = sub_box.text_frame
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = slide_data["subtitle"]
            sub_p.font.size = Pt(20)
            sub_p.font.color.rgb = colors["accent"]

        bullets = slide_data.get("bullets", [])
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(2.0) if i == 0 else Inches(1.8),
                Inches(12.0),
                Inches(5.0),
            )
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            for j, bullet in enumerate(bullets[:6]):
                bp = content_tf.paragraphs[0] if j == 0 else content_tf.add_paragraph()
                bp.text = f"  •  {bullet}"
                bp.font.size = Pt(18)
                bp.font.color.rgb = colors["text"]
                bp.space_after = Pt(8)

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    if not output_dir:
        output_dir = str(Path.home() / "Downloads")
    os.makedirs(output_dir, exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_topic = "".join(c for c in topic[:30] if c.isalnum() or c in " _-").strip().replace(" ", "_")
    filename = f"{safe_topic or 'pitch'}_{timestamp}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)

    try:
        if platform.system() == "Darwin":
            subprocess.run(["open", filepath], check=False, timeout=15)
        elif platform.system() == "Windows":
            os.startfile(filepath)  # type: ignore[attr-defined]
    except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
        pass

    return f"Pitch deck created: {filepath} ({len(slides_data)} slides). Opening now."


def _generate_slide_content(topic: str, slide_count: int, style: str) -> list[dict]:
    """Call Claude to generate structured slide content as JSON."""
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        return _fallback_slide_content(topic, slide_count)

    try:
        import anthropic

        client = anthropic.Anthropic(api_key=api_key)
        prompt = f"""Generate a {slide_count}-slide pitch deck about: {topic}

Return ONLY valid JSON, no markdown, no explanation. Format:
[
  {{
    "title": "Slide title",
    "subtitle": "Optional subtitle (only for slide 1)",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes for this slide"
  }}
]

Style: {style}. Make it compelling, concise, and professional.
Slide 1 should be the title slide. Last slide should be a call-to-action or summary."""

        response = client.messages.create(
            model="claude-haiku-4-5",
            max_tokens=2000,
            messages=[{"role": "user", "content": prompt}],
        )
        text = response.content[0].text.strip()
        if text.startswith("```"):
            text = text.split("```", 2)[1]
            if text.startswith("json"):
                text = text[4:]
            text = text.strip()
        parsed = json.loads(text)
        if isinstance(parsed, list) and parsed:
            return parsed[:slide_count]
    except Exception:  # noqa: BLE001
        pass
    return _fallback_slide_content(topic, slide_count)


def _fallback_slide_content(topic: str, slide_count: int) -> list[dict]:
    """Generate basic slide structure if Claude is unavailable."""
    slides = [
        {"title": topic, "subtitle": "Presentation", "bullets": [], "notes": "Title slide"},
        {"title": "Overview", "bullets": ["Key point 1", "Key point 2", "Key point 3"], "notes": ""},
        {"title": "The Problem", "bullets": ["Point 1", "Point 2", "Point 3"], "notes": ""},
        {"title": "Our Solution", "bullets": ["Solution 1", "Solution 2", "Solution 3"], "notes": ""},
        {"title": "Next Steps", "bullets": ["Action 1", "Action 2", "Action 3"], "notes": ""},
    ]
    return slides[:slide_count]
